"""Create final hackathon presentation — 15 slides covering all 10 scoring parameters."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x43, 0x61, 0xEE)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE6, 0x51, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox

def add_header(slide, title):
    add_bg(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), title, 32, True, WHITE)

def add_card(slide, left, top, width, height, title, desc, border_color=RGBColor(0xE0, 0xE0, 0xE0)):
    shape = add_shape(slide, left, top, width, height, WHITE)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
        title, 16, True, ACCENT_BLUE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(0.9),
        desc, 13, False, DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_BLUE)
add_shape(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "AI Across the Product\nDevelopment Lifecycle", 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
    "Intelligent Requirement-to-Test Traceability", 24, False, ACCENT_BLUE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
    "Hackathon 2026  |  AI Agents for PDLC", 18, False, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: The Problem  [Problem Definition]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "The Problem")
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), [
    "  Teams spend 4+ hours manually tracing requirements to test cases",
    "  Scattered documentation across formats (Word, PDF, Markdown, Text)",
    "  No automated traceability — coverage gaps go undetected until production",
    "  Test cases written before requirements are finalized → orphaned artifacts",
    "  Status reporting is manual, inconsistent, and time-consuming",
], 20, DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Our Approach  [Solution Architecture]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Our Approach")
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), [
    "  AI agents automate the entire requirement-to-test lifecycle",
    "  Upload documents in any format — PDF, Word, Markdown, Text",
    "  LLM-powered parsing extracts structured requirements",
    "  Bidirectional traceability: requirements→tests AND orphan detection",
    "  Multi-format compliance output: PDF, DOCX, TXT for auditors",
], 20, DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Architecture Overview  [Solution Architecture]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Architecture Overview")

boxes = [
    ("Frontend", "React + Vite\nUpload + Dashboard", Inches(0.5), Inches(2)),
    ("Backend", "Express + Prisma\nREST API", Inches(3.5), Inches(2)),
    ("LangGraph", "4-Agent Pipeline\nState Machine", Inches(6.5), Inches(2)),
    ("LLM + Embeddings", "Azure APIM\ngpt-5 + embeddings-3-large", Inches(9.5), Inches(2)),
]
for label, desc, left, top in boxes:
    shape = add_shape(slide, left, top, Inches(2.8), Inches(1.8), WHITE)
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.4), Inches(0.5),
        label, 16, True, ACCENT_BLUE, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.4), Inches(0.9),
        desc, 13, False, GRAY, PP_ALIGN.CENTER)

for x in [Inches(3.3), Inches(6.3), Inches(9.3)]:
    add_text_box(slide, x, Inches(2.6), Inches(0.4), Inches(0.5), "→", 24, True, ACCENT_BLUE, PP_ALIGN.CENTER)

bottom_boxes = [
    ("SQLite", "Prisma ORM\nRun + Document models", Inches(1.5), Inches(4.5)),
    ("Langfuse", "Observability\nLLM call tracing", Inches(5), Inches(4.5)),
    ("Puppeteer", "PDF Generation\nHTML → A4 PDF", Inches(8.5), Inches(4.5)),
]
for label, desc, left, top in bottom_boxes:
    shape = add_shape(slide, left, top, Inches(2.8), Inches(1.2), WHITE)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(2.4), Inches(0.4),
        label, 14, True, DARK_TEXT, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), Inches(2.4), Inches(0.6),
        desc, 12, False, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Detailed Architecture Diagram  [Solution Architecture]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Detailed System Architecture")

# User icon
shape = add_shape(slide, Inches(6), Inches(1.4), Inches(1.2), Inches(0.6), GRAY)
add_text_box(slide, Inches(6), Inches(1.42), Inches(1.2), Inches(0.5),
    "User", 11, True, WHITE, PP_ALIGN.CENTER)
# Arrow down
add_text_box(slide, Inches(6.4), Inches(1.95), Inches(0.4), Inches(0.4), "↓", 16, True, GRAY, PP_ALIGN.CENTER)

# Frontend
shape = add_shape(slide, Inches(4.2), Inches(2.4), Inches(4.8), Inches(0.7), ACCENT_BLUE)
add_text_box(slide, Inches(4.2), Inches(2.42), Inches(4.8), Inches(0.35),
    "Frontend — Vite + React", 12, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.2), Inches(2.72), Inches(4.8), Inches(0.35),
    "Home.tsx (Upload) · Dashboard.tsx (Results)", 9, False, WHITE, PP_ALIGN.CENTER)
# Arrow down
add_text_box(slide, Inches(6.4), Inches(3.1), Inches(0.4), Inches(0.4), "↓", 16, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.8), Inches(3.15), Inches(1.5), Inches(0.3),
    "POST /api/upload", 8, False, GRAY)

# Express API Layer
shape = add_shape(slide, Inches(4.2), Inches(3.5), Inches(4.8), Inches(0.65), GREEN)
add_text_box(slide, Inches(4.2), Inches(3.52), Inches(4.8), Inches(0.3),
    "Express API Layer", 12, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.2), Inches(3.8), Inches(4.8), Inches(0.3),
    "upload · analyze · results · report", 9, False, WHITE, PP_ALIGN.CENTER)

# Side: Parsers (left)
shape = add_shape(slide, Inches(1.5), Inches(3.5), Inches(1.8), Inches(0.65), RGBColor(0x5C, 0x6B, 0xC0))
add_text_box(slide, Inches(1.5), Inches(3.52), Inches(1.8), Inches(0.3),
    "Parsers", 10, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(1.8), Inches(0.3),
    "pdf · docx · md", 8, False, WHITE, PP_ALIGN.CENTER)
# Arrow: parsers ← API
add_text_box(slide, Inches(3.3), Inches(3.6), Inches(0.9), Inches(0.3), "← parse", 8, False, GRAY, PP_ALIGN.CENTER)

# Side: SQLite (right)
shape = add_shape(slide, Inches(10), Inches(3.5), Inches(2.2), Inches(0.65), ORANGE)
add_text_box(slide, Inches(10), Inches(3.52), Inches(2.2), Inches(0.3),
    "SQLite (Prisma)", 10, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(10), Inches(3.8), Inches(2.2), Inches(0.3),
    "Run + Document models", 8, False, WHITE, PP_ALIGN.CENTER)
# Arrow: API → SQLite
add_text_box(slide, Inches(9), Inches(3.6), Inches(1), Inches(0.3), "store →", 8, False, GRAY, PP_ALIGN.CENTER)

# Arrow down to pipeline
add_text_box(slide, Inches(6.4), Inches(4.15), Inches(0.4), Inches(0.4), "↓", 16, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.8), Inches(4.2), Inches(1.8), Inches(0.3),
    "invoke pipeline", 8, False, GRAY)

# LangGraph Pipeline box (dashed border)
pipeline_shape = add_shape(slide, Inches(2.8), Inches(4.6), Inches(7.4), Inches(2.3), WHITE)
pipeline_shape.line.color.rgb = GREEN
pipeline_shape.line.width = Pt(1.5)
pipeline_shape.line.dash_style = 2  # dash
add_text_box(slide, Inches(2.8), Inches(4.65), Inches(7.4), Inches(0.3),
    "LANGGRAPH PIPELINE", 9, True, GREEN, PP_ALIGN.CENTER)

# Agent 1: RequirementRefiner
shape = add_shape(slide, Inches(3), Inches(5.05), Inches(2.2), Inches(0.7), ACCENT_BLUE)
add_text_box(slide, Inches(3), Inches(5.07), Inches(2.2), Inches(0.35),
    "1. RequirementRefiner", 9, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(3), Inches(5.37), Inches(2.2), Inches(0.35),
    "Parse + embed requirements", 7, False, WHITE, PP_ALIGN.CENTER)

# Arrow 1→2
add_text_box(slide, Inches(5.2), Inches(5.2), Inches(0.4), Inches(0.3), "→", 14, True, GREEN, PP_ALIGN.CENTER)

# Agent 2: TestCaseGenerator
shape = add_shape(slide, Inches(5.6), Inches(5.05), Inches(2.2), Inches(0.7), GREEN)
add_text_box(slide, Inches(5.6), Inches(5.07), Inches(2.2), Inches(0.35),
    "2. TestCaseGenerator", 9, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.6), Inches(5.37), Inches(2.2), Inches(0.35),
    "LLM generates test cases", 7, False, WHITE, PP_ALIGN.CENTER)

# Arrow 2→3 (down)
add_text_box(slide, Inches(6.4), Inches(5.75), Inches(0.4), Inches(0.3), "↓", 14, True, GREEN, PP_ALIGN.CENTER)

# Agent 3: TraceabilityAgent
shape = add_shape(slide, Inches(5.6), Inches(6.05), Inches(2.2), Inches(0.7), RGBColor(0x5C, 0x6B, 0xC0))
add_text_box(slide, Inches(5.6), Inches(6.07), Inches(2.2), Inches(0.35),
    "3. TraceabilityAgent", 9, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.6), Inches(6.37), Inches(2.2), Inches(0.35),
    "Bidirectional matching", 7, False, WHITE, PP_ALIGN.CENTER)

# Arrow 3→4 (left)
add_text_box(slide, Inches(5.2), Inches(6.2), Inches(0.4), Inches(0.3), "←", 14, True, GREEN, PP_ALIGN.CENTER)

# Agent 4: ReportGenerator
shape = add_shape(slide, Inches(3), Inches(6.05), Inches(2.2), Inches(0.7), ORANGE)
add_text_box(slide, Inches(3), Inches(6.07), Inches(2.2), Inches(0.35),
    "4. ReportGenerator", 9, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(3), Inches(6.37), Inches(2.2), Inches(0.35),
    "HTML → Puppeteer → PDF", 7, False, WHITE, PP_ALIGN.CENTER)

# Side: LLM Client (left of pipeline)
shape = add_shape(slide, Inches(0.5), Inches(5.2), Inches(1.8), Inches(0.65), ORANGE)
add_text_box(slide, Inches(0.5), Inches(5.22), Inches(1.8), Inches(0.3),
    "LLM Client", 9, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(5.5), Inches(1.8), Inches(0.3),
    "Azure APIM", 7, False, WHITE, PP_ALIGN.CENTER)
# Arrow: Agent1 → LLM
add_text_box(slide, Inches(2.3), Inches(5.3), Inches(0.7), Inches(0.3), "← call", 7, False, GRAY, PP_ALIGN.CENTER)

# Side: Embeddings API (right of pipeline)
shape = add_shape(slide, Inches(10.5), Inches(6.1), Inches(2.2), Inches(0.65), RGBColor(0x5C, 0x6B, 0xC0))
add_text_box(slide, Inches(10.5), Inches(6.12), Inches(2.2), Inches(0.3),
    "Embeddings API", 9, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.5), Inches(6.4), Inches(2.2), Inches(0.3),
    "text-embedding-3-large", 7, False, WHITE, PP_ALIGN.CENTER)
# Arrow: Agent3 → Embeddings
add_text_box(slide, Inches(7.8), Inches(6.2), Inches(2.7), Inches(0.3), "embed →", 7, False, GRAY, PP_ALIGN.CENTER)

# Langfuse — SOLID, REQUIRED (not faded)
shape = add_shape(slide, Inches(10.5), Inches(5.2), Inches(2.2), Inches(0.65), RGBColor(0x7C, 0x4D, 0xFF))
shape.line.color.rgb = RGBColor(0x7C, 0x4D, 0xFF)
shape.line.width = Pt(2)
add_text_box(slide, Inches(10.5), Inches(5.22), Inches(2.2), Inches(0.3),
    "Langfuse", 10, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.5), Inches(5.5), Inches(2.2), Inches(0.3),
    "Observability (Required)", 7, False, WHITE, PP_ALIGN.CENTER)
# Dashed line from Langfuse to pipeline
add_text_box(slide, Inches(9.8), Inches(5.4), Inches(0.7), Inches(0.3), "← trace", 7, False, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: The 4 AI Agents  [AI Implementation]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "The 4 AI Agents")

agents = [
    ("1. Requirement Refiner", "Parses uploaded documents using LLM.\nExtracts structured requirements with IDs.\nComputes embeddings for semantic matching."),
    ("2. Test Case Generator", "Generates test cases at 4 levels:\nUnit, Integration, System, Acceptance.\nFalls back to keyword-based generation."),
    ("3. Traceability Agent", "Bidirectional matching:\nForward: requirements → test cases\nReverse: orphan test case detection"),
    ("4. Report Generator", "Renders HTML traceability report.\nConverts to PDF via Puppeteer.\nFalls back to HTML if Puppeteer fails."),
]
for i, (title, desc) in enumerate(agents):
    left = Inches(0.5) + (i * Inches(3.15))
    add_card(slide, left, Inches(1.8), Inches(2.9), Inches(4.5), title, desc, ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Semantic Similarity  [AI Implementation]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Semantic Similarity: How It Works")

steps = [
    ("Text Input", "\"Users must log in\nwith Google OAuth\""),
    ("Embedding", "text-embedding-3-large\n→ 1536-dimension vector"),
    ("Comparison", "Cosine similarity\nbetween vectors"),
    ("Match", "85% similarity\n= Match found"),
]
for i, (label, desc) in enumerate(steps):
    left = Inches(0.8) + (i * Inches(3.1))
    add_card(slide, left, Inches(2), Inches(2.5), Inches(2), label, desc, ACCENT_BLUE)

for x in [Inches(3.3), Inches(6.4), Inches(9.5)]:
    add_text_box(slide, x, Inches(2.7), Inches(0.4), Inches(0.5), "→", 24, True, ACCENT_BLUE, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
    "Fallback: When embeddings API is unavailable, degrades to word-level Jaccard similarity\nwith stem matching and synonym dictionary — never blocks the user.",
    14, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.8),
    "\"Google Login\" vs \"Verify OAuth sign-in\"  →  Semantic: 85%  |  Keyword: 15%",
    16, True, DARK_TEXT, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: AI Appropriateness  [AI Appropriateness] — NEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "AI Appropriateness")

# Left column — AI Handles
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
    "AI Handles", 22, True, ACCENT_BLUE, PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5), [
    "  Unstructured → structured requirement extraction",
    "  Semantic test case generation (4 testing levels)",
    "  Embedding-based traceability matching",
    "  Domain-adaptive prompt engineering",
], 16, DARK_TEXT)

# Right column — Traditional Code Handles
add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.6),
    "Traditional Code Handles", 22, True, GRAY, PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(3.5), [
    "  File parsing (PDF, DOCX, Markdown, TXT)",
    "  Report rendering (HTML → PDF via Puppeteer)",
    "  Database operations (Prisma + SQLite)",
    "  Cosine similarity computation (deterministic)",
], 16, DARK_TEXT)

# Divider
add_shape(slide, Inches(6.5), Inches(2.3), Inches(0.04), Inches(3), RGBColor(0xE0, 0xE0, 0xE0))

# Callout
shape = add_shape(slide, Inches(2), Inches(5.8), Inches(9), Inches(1), WHITE)
shape.line.color.rgb = ACCENT_BLUE
shape.line.width = Pt(2)
add_text_box(slide, Inches(2.3), Inches(5.95), Inches(8.4), Inches(0.7),
    "AI is used where it genuinely adds value — not everywhere", 18, True, ACCENT_BLUE, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Demo Walkthrough  [Working Demo]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Demo Walkthrough")

flow_steps = [
    ("1. Upload", "Drag & drop\nrequirements +\ntest cases"),
    ("2. Analyze", "AI agents parse,\ngenerate tests,\nfind gaps"),
    ("3. Dashboard", "Coverage stats,\ntraceability matrix,\norphan detection"),
    ("4. PDF Report", "Professional report\nwith full analysis\nfor stakeholders"),
]
for i, (label, desc) in enumerate(flow_steps):
    left = Inches(0.8) + (i * Inches(3.1))
    shape = add_shape(slide, left + Inches(0.8), Inches(1.8), Inches(0.8), Inches(0.8), ACCENT_BLUE)
    add_text_box(slide, left + Inches(0.8), Inches(1.85), Inches(0.8), Inches(0.7),
        str(i+1), 24, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.8), Inches(2.5), Inches(0.5),
        label, 18, True, DARK_TEXT, PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(3.3), Inches(2.5), Inches(1.5),
        desc, 14, False, GRAY, PP_ALIGN.CENTER)

for x in [Inches(3.6), Inches(6.7), Inches(9.8)]:
    add_text_box(slide, x, Inches(2), Inches(0.4), Inches(0.5), "→", 24, True, ACCENT_BLUE, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Tech Stack  [Engineering Quality]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Tech Stack")

stack = [
    ("Frontend", "React 19 + Vite 6\nTypeScript, React Router"),
    ("Backend", "Express.js + Node.js\nTypeScript (strict mode)"),
    ("AI Orchestration", "LangGraph.js\n4-agent state machine"),
    ("LLM", "Azure APIM\ngpt-5.2 + embeddings-3-large"),
    ("Database", "SQLite + Prisma ORM\nZero-install, schema-first"),
    ("PDF Generation", "Puppeteer\nHTML → A4 PDF"),
    ("Observability", "Langfuse\nEvery LLM call traced"),
    ("File Parsing", "pdf-parse, mammoth\nPDF, Word, Markdown, Text"),
]
for i, (label, desc) in enumerate(stack):
    col = i % 4
    row = i // 4
    left = Inches(0.5) + (col * Inches(3.15))
    top = Inches(1.6) + (row * Inches(2.5))
    add_card(slide, left, top, Inches(2.9), Inches(2), label, desc)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: AI Impact  [AI Impact]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "AI Impact: Before vs After")

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
    "Before (Manual)", 22, True, RED, PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4), [
    "  4+ hours per traceability analysis",
    "  Manual requirement-to-test matching",
    "  Keyword search misses semantic links",
    "  Coverage gaps discovered in production",
    "  No observability into analysis process",
], 16, DARK_TEXT)

add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.6),
    "After (AI-Powered)", 22, True, GREEN, PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4), [
    "  30 seconds end-to-end analysis",
    "  Automated semantic matching via embeddings",
    "  Catches \"Google Login\" = \"OAuth sign-in\"",
    "  Gaps flagged before code is written",
    "  Langfuse traces every AI decision",
], 16, DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Engineering Quality  [Engineering Quality]
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Engineering Quality")

qualities = [
    ("Graceful Fallbacks", "Every AI component has a non-AI fallback.\nLLM fails → regex parsing.\nEmbeddings fail → word-level matching.\nPuppeteer fails → HTML report."),
    ("Type Safety", "TypeScript strict mode throughout.\nShared AgentState types across all agents.\nPrisma schema-first database design."),
    ("Observability", "Langfuse integration traces every LLM call.\nInput/output logged for debugging.\nGraceful degradation when not configured."),
    ("Clean Architecture", "Separation: agents, routes, parsers.\nMonorepo with npm workspaces.\nZero-install database (SQLite)."),
]
for i, (title, desc) in enumerate(qualities):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + (col * Inches(6.3))
    top = Inches(1.6) + (row * Inches(2.7))
    add_card(slide, left, top, Inches(5.8), Inches(2.3), title, desc, ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Business Relevance & Viability  [Business Relevance + Benefits] — NEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Business Relevance & Viability")

biz_cards = [
    ("Multi-Domain", "Same engine works for:\nSoftware, EV, Civil Engineering,\nPayments, Aerospace.\nSwap the template, not the code."),
    ("Compliance Ready", "PDF/TXT/DOCX output for\nauditors in regulated industries.\nAutomotive (ISO 26262),\nBanking (PCI-DSS)."),
    ("Zero Setup", "SQLite + local LLM calls.\nNo vector DB, no cloud\ndependencies required.\nWorks offline with fallbacks."),
    ("Graceful Degradation", "Never blocks the user.\nLLM unavailable → regex parsing.\nEmbeddings fail → word matching.\nPuppeteer fails → HTML report."),
]
for i, (title, desc) in enumerate(biz_cards):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + (col * Inches(6.3))
    top = Inches(1.6) + (row * Inches(2.7))
    add_card(slide, left, top, Inches(5.8), Inches(2.3), title, desc)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: What Sets Us Apart  [Differentiation] — NEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "What Sets Us Apart")

diff_cards = [
    ("Bidirectional\nTraceability", "Most tools only map\nrequirements → tests.\n\nWe also detect orphan\ntest cases that have\nno linked requirement."),
    ("Domain-Agnostic\nEngine", "Same pipeline works for\nsoftware, automotive,\ncivil engineering,\nbanking.\n\n5 industry samples included."),
    ("Backwards\nCompatible", "Handles real-world chaos:\ntest cases without\nrequirements,\npartial uploads,\nmissing documents."),
]
for i, (title, desc) in enumerate(diff_cards):
    left = Inches(0.5) + (i * Inches(4.2))
    add_card(slide, left, Inches(1.8), Inches(3.9), Inches(4.5), title, desc, ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14: Future Enhancements
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Future Enhancements")

enhancements = [
    ("1. RAG + Vector DB", "ChromaDB/Pinecone for large-scale\nsemantic document search."),
    ("2. Design Review\nSummarization", "AI agent that summarizes design\ndocuments and links to requirements."),
    ("3. Multi-Run History", "Track coverage trends across\nsprint cycles and versions."),
    ("4. Chat Q&A Interface", "Natural language queries:\n\"Why is Req-7 not covered?\""),
    ("5. AI Assumptions\nClarification", "AI surfaces its assumptions before\ngenerating test cases — user\nclarifies upfront, saving tokens."),
]
# Row 1: 3 cards
for i in range(3):
    left = Inches(0.5) + (i * Inches(4.2))
    add_card(slide, left, Inches(1.6), Inches(3.9), Inches(2.3), enhancements[i][0], enhancements[i][1])
# Row 2: 2 cards centered
for i in range(2):
    left = Inches(2.6) + (i * Inches(4.2))
    add_card(slide, left, Inches(4.2), Inches(3.9), Inches(2.3), enhancements[i+3][0], enhancements[i+3][1], ORANGE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15: Q&A
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_BLUE)
add_shape(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1),
    "Thank You", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
    "Questions?", 32, False, ACCENT_BLUE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
    "AI Across the Product Development Lifecycle  |  Hackathon 2026", 16, False, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "..", "docs", "Architecture", "PDLC-Final-Presentation-v2.pptx")
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
