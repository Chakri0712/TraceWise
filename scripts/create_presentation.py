"""Create hackathon judge presentation for PDLC AI Agents project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x43, 0x61, 0xEE)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return txBox

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(7.5), DARK_BLUE)
add_shape(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "AI Across the Product\nDevelopment Lifecycle", 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
    "Intelligent Requirement-to-Test Traceability", 24, False, ACCENT_BLUE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
    "Hackathon 2026  |  AI Agents for PDLC", 18, False, GRAY, PP_ALIGN.CENTER)

# ── Slide 2: The Problem ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "The Problem", 32, True, WHITE)

problems = [
    "Manual handoffs between requirements, design, and testing teams",
    "Scattered documentation across formats (Word, PDF, Markdown)",
    "No automated traceability from requirement to test case",
    "Coverage gaps go undetected until production issues arise",
    "Status reporting is manual, inconsistent, and time-consuming",
]
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
    [f"  {p}" for p in problems], 20, DARK_TEXT)

# ── Slide 3: Our Approach ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "Our Approach", 32, True, WHITE)

approach = [
    "AI agents that automate the entire requirement-to-test lifecycle",
    "Upload documents in any format (PDF, Word, Markdown, Text)",
    "LLM-powered parsing extracts structured requirements",
    "Semantic similarity matching finds test-to-requirement links",
    "Automated gap detection with AI-generated test suggestions",
    "Professional PDF reports with traceability matrix",
]
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
    [f"  {a}" for a in approach], 20, DARK_TEXT)

# ── Slide 4: Architecture Overview ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "Architecture Overview", 32, True, WHITE)

# Architecture boxes
boxes = [
    ("Frontend", "React + Vite\nUpload + Dashboard", Inches(0.5), Inches(2)),
    ("Backend", "Express + Prisma\nREST API", Inches(3.5), Inches(2)),
    ("LangGraph", "4-Agent Pipeline\nState Machine", Inches(6.5), Inches(2)),
    ("LLM + Embeddings", "Azure APIM\ngpt-5 + embeddings-3-large", Inches(9.5), Inches(2)),
]
for label, desc, left, top in boxes:
    shape = add_shape(slide, left, top, Inches(2.8), Inches(1.8), WHITE)
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.4), Inches(0.5),
        label, 16, True, ACCENT_BLUE, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.4), Inches(0.9),
        desc, 13, False, GRAY, PP_ALIGN.CENTER)

# Arrows
for x in [Inches(3.3), Inches(6.3), Inches(9.3)]:
    add_text_box(slide, x, Inches(2.6), Inches(0.4), Inches(0.5),
        "→", 24, True, ACCENT_BLUE, PP_ALIGN.CENTER)

# Bottom boxes
bottom_boxes = [
    ("SQLite", "Prisma ORM\nRun + Document models", Inches(1.5), Inches(4.5)),
    ("Langfuse", "Observability\nAgent tracing", Inches(5), Inches(4.5)),
    ("Puppeteer", "PDF Generation\nHTML → PDF", Inches(8.5), Inches(4.5)),
]
for label, desc, left, top in bottom_boxes:
    shape = add_shape(slide, left, top, Inches(2.8), Inches(1.2), WHITE)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(2.4), Inches(0.4),
        label, 14, True, DARK_TEXT, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), Inches(2.4), Inches(0.6),
        desc, 12, False, GRAY, PP_ALIGN.CENTER)

# ── Slide 5: The 4 AI Agents ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "The 4 AI Agents", 32, True, WHITE)

agents = [
    ("1. Requirement Refiner", "Parses uploaded documents using LLM.\nExtracts structured requirements with IDs and priorities.\nComputes embeddings for semantic matching."),
    ("2. Test Case Generator", "Generates detailed test cases from requirements using LLM.\nCreates 1-2 tests per requirement.\nFalls back to keyword-based generation if LLM fails."),
    ("3. Traceability Agent", "Matches uploaded test cases to requirements.\nUses cosine similarity on embeddings (threshold: 0.6).\nIdentifies coverage gaps."),
    ("4. Report Generator", "Renders HTML report with traceability matrix.\nConverts to PDF via Puppeteer.\nFalls back to HTML if Puppeteer fails."),
]

for i, (title, desc) in enumerate(agents):
    left = Inches(0.5) + (i * Inches(3.15))
    shape = add_shape(slide, left, Inches(1.8), Inches(2.9), Inches(4.5), WHITE)
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.2), Inches(2), Inches(2.5), Inches(0.6),
        title, 15, True, ACCENT_BLUE, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(2.7), Inches(2.5), Inches(3.3),
        desc, 12, False, DARK_TEXT)

# ── Slide 6: Semantic Similarity ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "Semantic Similarity: How It Works", 32, True, WHITE)

steps = [
    ("Text Input", "\"Users must log in\nwith Google OAuth\""),
    ("Embedding", "text-embedding-3-large\n→ 1536-dimension vector"),
    ("Comparison", "Cosine similarity\nbetween vectors"),
    ("Match", "63% similarity\n= Match found"),
]

for i, (label, desc) in enumerate(steps):
    left = Inches(0.8) + (i * Inches(3.1))
    shape = add_shape(slide, left, Inches(2), Inches(2.5), Inches(2), WHITE)
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.2), Inches(2.2), Inches(2.1), Inches(0.5),
        label, 16, True, ACCENT_BLUE, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(2.7), Inches(2.1), Inches(1.2),
        desc, 13, False, DARK_TEXT, PP_ALIGN.CENTER)

# Arrow between boxes
for x in [Inches(3.3), Inches(6.4), Inches(9.5)]:
    add_text_box(slide, x, Inches(2.7), Inches(0.4), Inches(0.5),
        "→", 24, True, ACCENT_BLUE, PP_ALIGN.CENTER)

# Fallback note
add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(1),
    "Fallback: If embeddings API is unavailable, the system degrades gracefully to\nword-level Jaccard similarity with stem matching and synonym dictionary.",
    14, False, GRAY, PP_ALIGN.CENTER)

# Comparison example
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(1),
    "Example: \"Google Login\" vs \"Verify OAuth sign-in works\"  →  Semantic match (0.85)  |  Keyword match (0.15)",
    16, True, DARK_TEXT, PP_ALIGN.CENTER)

# ── Slide 7: Demo Walkthrough ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "Demo Walkthrough", 32, True, WHITE)

flow_steps = [
    ("1. Upload", "Drag & drop\nrequirements +\ntest cases"),
    ("2. Analyze", "AI agents parse,\ngenerate tests,\nfind gaps"),
    ("3. Dashboard", "Coverage stats,\ntraceability matrix,\ngap details"),
    ("4. PDF Report", "Professional report\nwith full analysis\nfor stakeholders"),
]

for i, (label, desc) in enumerate(flow_steps):
    left = Inches(0.8) + (i * Inches(3.1))
    # Number circle
    shape = add_shape(slide, left + Inches(0.8), Inches(1.8), Inches(0.8), Inches(0.8), ACCENT_BLUE)
    add_text_box(slide, left + Inches(0.8), Inches(1.85), Inches(0.8), Inches(0.7),
        str(i+1), 24, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.8), Inches(2.5), Inches(0.5),
        label, 18, True, DARK_TEXT, PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(3.3), Inches(2.5), Inches(1.5),
        desc, 14, False, GRAY, PP_ALIGN.CENTER)

# Arrow between steps
for x in [Inches(3.6), Inches(6.7), Inches(9.8)]:
    add_text_box(slide, x, Inches(2), Inches(0.4), Inches(0.5),
        "→", 24, True, ACCENT_BLUE, PP_ALIGN.CENTER)

# ── Slide 8: Tech Stack ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "Tech Stack", 32, True, WHITE)

stack = [
    ("Frontend", "React 18 + Vite 6\nTypeScript, React Router"),
    ("Backend", "Express.js + Node.js\nTypeScript (strict mode)"),
    ("AI Orchestration", "LangGraph.js\n4-agent state machine"),
    ("LLM", "Azure APIM\ngpt-5-mini + embeddings-3-large"),
    ("Database", "SQLite + Prisma ORM\nZero-install, schema-first"),
    ("PDF Generation", "Puppeteer\nHTML → PDF (A4 format)"),
    ("Observability", "Langfuse\nAgent tracing (optional)"),
    ("File Parsing", "pdf-parse, mammoth, marked\nPDF, Word, Markdown, Text"),
]

for i, (label, desc) in enumerate(stack):
    col = i % 4
    row = i // 4
    left = Inches(0.5) + (col * Inches(3.15))
    top = Inches(1.6) + (row * Inches(2.5))
    shape = add_shape(slide, left, top, Inches(2.9), Inches(2), WHITE)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.5), Inches(0.5),
        label, 15, True, ACCENT_BLUE, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.5), Inches(1.1),
        desc, 13, False, DARK_TEXT, PP_ALIGN.LEFT)

# ── Slide 9: AI Impact ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "AI Impact: Before vs After", 32, True, WHITE)

# Before column
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
    "Before (Manual)", 22, True, RED, PP_ALIGN.CENTER)
before_items = [
    "  Manual requirement-to-test matching",
    "  Keyword-based search misses semantic links",
    "  Coverage gaps discovered in production",
    "  Hours spent on traceability reports",
    "  No observability into analysis process",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4),
    before_items, 16, DARK_TEXT)

# After column
add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.6),
    "After (AI-Powered)", 22, True, GREEN, PP_ALIGN.CENTER)
after_items = [
    "  Automated semantic matching via embeddings",
    "  Catches \"Google Login\" = \"OAuth sign-in\"",
    "  Gaps flagged before code is written",
    "  Seconds to generate full traceability report",
    "  Langfuse traces every AI decision",
]
add_bullet_list(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4),
    after_items, 16, DARK_TEXT)

# ── Slide 10: Engineering Quality ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "Engineering Quality", 32, True, WHITE)

qualities = [
    ("Graceful Fallbacks", "Every AI component has a non-AI fallback.\nLLM fails → regex parsing.\nEmbeddings fail → word-level matching.\nPuppeteer fails → HTML report."),
    ("Type Safety", "TypeScript strict mode throughout.\nShared AgentState types across all agents.\nPrisma schema-first database design."),
    ("Observability", "Langfuse integration traces every LLM call.\nInput/output logged for debugging.\nGraceful degradation when not configured."),
    ("Clean Architecture", "Separation of concerns: agents, routes, parsers.\nMonorepo with npm workspaces.\nZero-install database (SQLite)."),
]

for i, (title, desc) in enumerate(qualities):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + (col * Inches(6.3))
    top = Inches(1.6) + (row * Inches(2.7))
    shape = add_shape(slide, left, top, Inches(5.8), Inches(2.3), WHITE)
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(5.2), Inches(0.5),
        title, 18, True, ACCENT_BLUE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), Inches(5.2), Inches(1.4),
        desc, 14, False, DARK_TEXT)

# ── Slide 11: Future Enhancements ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
    "Future Enhancements", 32, True, WHITE)

enhancements = [
    ("RAG + Vector Database", "ChromaDB/Pinecone for large-scale document storage.\nChunk requirements into searchable vectors.\nRetrieve relevant context for each analysis."),
    ("Design Review Summarization", "AI agent that summarizes design documents.\nExtracts key decisions and action items.\nLinks design decisions to requirements."),
    ("Multi-Run History", "Track analysis runs over time.\nCompare coverage across versions.\nTrend dashboards for stakeholder reporting."),
    ("Chat Q&A Interface", "RAG-powered chat on requirement gaps.\n\"Why is Req-7 not covered?\" → AI explanation.\nInteractive gap exploration."),
]

for i, (title, desc) in enumerate(enhancements):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + (col * Inches(6.3))
    top = Inches(1.6) + (row * Inches(2.7))
    shape = add_shape(slide, left, top, Inches(5.8), Inches(2.3), WHITE)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(5.2), Inches(0.5),
        title, 18, True, ACCENT_BLUE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), Inches(5.2), Inches(1.4),
        desc, 14, False, DARK_TEXT)

# ── Slide 12: Q&A ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(7.5), DARK_BLUE)
add_shape(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1),
    "Thank You", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
    "Questions?", 32, False, ACCENT_BLUE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
    "AI Across the Product Development Lifecycle  |  Hackathon 2026", 16, False, GRAY, PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(__file__), "docs", "PDLC-Architecture-Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to {output_path}")
